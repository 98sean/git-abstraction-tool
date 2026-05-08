"""Generate presentation.pptx for the Honda Team project."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RED   = RGBColor(0xCC, 0x00, 0x00)
DARK  = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
CREAM = RGBColor(0xED, 0xE8, 0xE0)
GRAY  = RGBColor(0x88, 0x88, 0x88)
MID   = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def bg(sl, color):
    sh = sl.shapes.add_shape(1, 0, 0, W, H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.line.width = 0


def box(sl, x, y, w, h, fill=None):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.line.fill.background()
    sh.line.width = 0
    return sh


def tb(sl, x, y, w, h, text, size=13, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT):
    if color is None:
        color = MID
    t = sl.shapes.add_textbox(x, y, w, h)
    t.word_wrap = True
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return t


def lbl(sl, x, y, w, text, color=GRAY):
    tb(sl, x, y, w, Inches(0.3), text.upper(),
       size=8, bold=True, color=color)


def h1(sl, x, y, w, text, size=36, color=DARK, align=PP_ALIGN.LEFT):
    tb(sl, x, y, w, Inches(1.4), text, size=size, bold=True,
       color=color, align=align)


def rule(sl, x, y, w):
    sh = sl.shapes.add_shape(1, x, y, w, Pt(2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    sh.line.fill.background()
    sh.line.width = 0


def card(sl, x, y, w, h, head, body, hsize=12, bsize=11,
         accent=RED, bg_color=WHITE):
    box(sl, x, y, w, h, fill=bg_color)
    bar = box(sl, x, y, Inches(0.07), h, fill=accent)
    if head:
        tb(sl, x+Inches(0.17), y+Inches(0.1), w-Inches(0.25),
           Inches(0.33), head, size=hsize, bold=True, color=DARK)
    tb(sl, x+Inches(0.17),
       y+(Inches(0.1) if not head else Inches(0.42)),
       w-Inches(0.25),
       h-(Inches(0.1) if not head else Inches(0.5)),
       body, size=bsize, color=MID)


def bullets(sl, x, y, w, items, size=12, gap=Inches(0.46),
            dot_col=RED, text_col=MID):
    cy = y
    for item in items:
        dot = sl.shapes.add_shape(9, x, cy+Inches(0.1),
                                  Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_col
        dot.line.fill.background()
        dot.line.width = 0
        tb(sl, x+Inches(0.21), cy, w-Inches(0.21), gap,
           item, size=size, color=text_col)
        cy += gap


def pill(sl, x, y, text):
    bx = box(sl, x, y, Inches(2.1), Inches(0.3), fill=RED)
    tb(sl, x+Inches(0.1), y+Inches(0.04), Inches(1.9), Inches(0.25),
       text.upper(), size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def quote(sl, x, y, w, text, size=12):
    box(sl, x, y, Inches(0.06), Inches(0.68), fill=RED)
    box(sl, x+Inches(0.06), y, w-Inches(0.06), Inches(0.68),
        fill=RGBColor(0xFD, 0xF0, 0xF0))
    tb(sl, x+Inches(0.2), y+Inches(0.08), w-Inches(0.28),
       Inches(0.55), '"' + text + '"',
       size=size, italic=True, color=RGBColor(0x55, 0x22, 0x22))


# ============================================================
# SLIDE 1 -- COVER
# ============================================================
sl = new_slide()
bg(sl, DARK)
lbl(sl, Inches(1), Inches(1.1), Inches(11),
    "Honda x 99P Labs - Spring 2026",
    color=RGBColor(0x77, 0x77, 0x77))
tb(sl, Inches(1), Inches(1.45), Inches(11.3), Inches(2.5),
   "Code Repo Knowledge Layer\n-> Git Abstraction Tool",
   size=40, bold=True, color=WHITE)
tb(sl, Inches(1), Inches(4.1), Inches(9.5), Inches(1.2),
   "A semester-long journey from developer onboarding tooling to a desktop app\n"
   "that makes version control accessible to everyone.",
   size=15, color=RGBColor(0xAA, 0xAA, 0xAA))

# ============================================================
# SLIDE 2 -- SECTION: Part 1
# ============================================================
sl = new_slide()
bg(sl, RED)
lbl(sl, Inches(1), Inches(2.7), Inches(10), "Part 1",
    color=RGBColor(0xFF, 0xCC, 0xCC))
h1(sl, Inches(1), Inches(3.05), Inches(11.3),
   "Code Repo Knowledge Layer", size=38, color=WHITE)
tb(sl, Inches(1), Inches(4.55), Inches(9), Inches(0.6),
   "Our initial project -- Repo Onboarding Buddy",
   size=17, color=RGBColor(0xFF, 0xDD, 0xDD))

# ============================================================
# SLIDE 3 -- The Problem
# ============================================================
sl = new_slide()
bg(sl, LIGHT)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "Code Repo Knowledge Layer -- Problem")
h1(sl, Inches(0.6), Inches(0.62), Inches(12),
   "The Problem We Started With", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

bullets(sl, Inches(0.6), Inches(1.6), Inches(5.85),
    [
        "Explicit knowledge -- code, comments, README. Already in the repo, "
        "but hard to navigate at scale.",
        "Structural knowledge -- 'if you change this, that breaks.' Implicit "
        "in code but hard to trace manually.",
        "Contextual knowledge -- 'why was this written this way?' Lives only "
        "in people's heads. Lost completely when team communication is absent.",
    ], size=12, gap=Inches(0.82))

card(sl, Inches(7.0), Inches(1.6), Inches(5.75), Inches(1.1),
     "Root Cause",
     "Teams don't communicate around code -- so new members must read "
     "everything manually, and it takes too long.", bsize=11)
card(sl, Inches(7.0), Inches(2.85), Inches(5.75), Inches(1.1),
     "Gap in existing tools",
     "GitHub Copilot and search tools handle explicit knowledge only. "
     "Structural and contextual knowledge still require asking people.", bsize=11)
card(sl, Inches(7.0), Inches(4.1), Inches(5.75), Inches(1.0),
     "Our goal",
     "Build a knowledge layer covering all three types without leaving "
     "the editor.", bsize=11)

# ============================================================
# SLIDE 4 -- Five Directions
# ============================================================
sl = new_slide()
bg(sl, CREAM)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "Code Repo Knowledge Layer -- Brainstorm")
h1(sl, Inches(0.6), Inches(0.62), Inches(12),
   "Five Directions We Explored", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

left_items = [
    ('1. AI Coaching Agent -- "Repo Onboarding Buddy"',
     "Ask 'how does the data pipeline work?' and get an answer grounded in "
     "code structure, dependencies, and git history -- not just docs."),
    ("2. Auto-Documentation",
     "Living docs that auto-update when code changes. Four layers: one-liner, "
     "feature description, architecture context, business impact."),
    ("3. Usage Tracking Layer",
     "Git history + AST analysis -> Core Assets, Stable components, Dead Code "
     "candidates. Bus factor alerts. Tableau treemap visualization."),
]
right_items = [
    ("4. Interactive Code Map",
     "Navigable graph of the entire repo. Click to drill down. Data lineage "
     "and cross-team dependency discovery."),
    ("5. Communication Gap Filler",
     "Weekly AI newsletter of what changed. Cross-team impact alerts. "
     "'Ask the repo' Slack bot."),
]

y = Inches(1.6)
for head, body in left_items:
    card(sl, Inches(0.6), y, Inches(5.9), Inches(0.97),
         head, body, hsize=11, bsize=10)
    y += Inches(1.07)

y = Inches(1.6)
for head, body in right_items:
    card(sl, Inches(6.9), y, Inches(5.9), Inches(0.97),
         head, body, hsize=11, bsize=10)
    y += Inches(1.07)

card(sl, Inches(6.9), Inches(3.75), Inches(5.9), Inches(0.85),
     "One-line product definition",
     '"A VS Code Extension that lets you understand every context of a repo '
     'without leaving the editor."',
     hsize=10, bsize=10, accent=GRAY)

# ============================================================
# SLIDE 5 -- What We Built
# ============================================================
sl = new_slide()
bg(sl, WHITE)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "Code Repo Knowledge Layer -- Phase 1 Deliverables")
h1(sl, Inches(0.6), Inches(0.62), Inches(12),
   "What We Actually Built", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.58), Inches(5.9), Inches(0.35),
   "Repo Onboarding Buddy (Chanyoung, Seungyeon, Jaeyoon)",
   size=12, bold=True, color=RED)
bullets(sl, Inches(0.6), Inches(2.0), Inches(5.9),
    ["VS Code Extension UI with sidebar chat panel",
     "Context-aware Q&A: agent knows the currently open file",
     "AST-based code parsing with tree-sitter (no AI for structure)",
     "Dependency graph via NetworkX; impact analysis on right-click",
     "Hover metadata: one-liner, caller count, bus factor warning"],
    size=11, gap=Inches(0.44))

tb(sl, Inches(6.9), Inches(1.58), Inches(5.9), Inches(0.35),
   "Usage Tracking Layer (Jonghoon, Gwanho, Yunha)",
   size=12, bold=True, color=RED)
bullets(sl, Inches(6.9), Inches(2.0), Inches(5.9),
    ["CLI tool 'repoknow scan' + FastAPI backend",
     "PyDriller for git history extraction (churn, ownership, bus factor)",
     "CodeQL for static import graph -> internal_edges.csv",
     "Scoring: Core Assets / Stable / Dead Candidates / NonCode",
     "REST API: /impact, /bus-factor-alerts, /top-imports"],
    size=11, gap=Inches(0.44))

# ============================================================
# SLIDE 6 -- SECTION: Part 2
# ============================================================
sl = new_slide()
bg(sl, RED)
lbl(sl, Inches(1), Inches(2.7), Inches(10), "Part 2",
    color=RGBColor(0xFF, 0xCC, 0xCC))
h1(sl, Inches(1), Inches(3.05), Inches(11.3),
   "The Pivot", size=38, color=WHITE)
tb(sl, Inches(1), Inches(4.2), Inches(9), Inches(0.6),
   "Why we shifted from Repo Buddy to Git Abstraction Tool",
   size=17, color=RGBColor(0xFF, 0xDD, 0xDD))

# ============================================================
# SLIDE 7 -- Pivot: Barriers
# ============================================================
sl = new_slide()
bg(sl, LIGHT)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10), "The Pivot -- What We Hit")
h1(sl, Inches(0.6), Inches(0.62), Inches(12),
   "Barriers We Encountered", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

bullets(sl, Inches(0.6), Inches(1.6), Inches(5.9),
    ["No shared DB or server -- Ryan's guidance: 'don't spend any money.' "
     "SQLite local-only was the ceiling for Phase 1.",
     "Cold start problem -- Q&A bot needs seniors to answer first, but they "
     "already use Copilot for code explanation.",
     "Weak value proposition -- 'What does this function do?' is already "
     "answered by Claude and Copilot. Not differentiated.",
     "Tribal knowledge gap -- 'Why was this written this way?' requires git "
     "history + PR discussion, hard to build reliably."],
    size=12, gap=Inches(0.75))

card(sl, Inches(7.0), Inches(1.6), Inches(5.8), Inches(1.25),
     "Ryan's Feedback -- 3/6 Honda Meeting",
     '"Think of this project as a new approach to version control. '
     'A new generation of GitHub or Git."', bsize=12)
card(sl, Inches(7.0), Inches(3.0), Inches(5.8), Inches(1.5),
     "Ryan's Feedback -- 4/3 Honda Meeting",
     '"How do we make vibe coding easier? Ppl do not know what they want -- '
     'AI can give a high-level plan. Think about the vibe coders using phones '
     'or tablets to code."', bsize=12)
card(sl, Inches(7.0), Inches(4.65), Inches(5.8), Inches(0.9),
     "New Vision",
     "A next-gen git tool for non-technical users -- local-first, "
     "no server cost, AI-optional.", bsize=12, accent=GRAY)

# ============================================================
# SLIDE 8 -- Shift Diagram
# ============================================================
sl = new_slide()
bg(sl, CREAM)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "The Pivot -- Before vs. After")
h1(sl, Inches(0.6), Inches(0.62), Inches(12), "What Changed", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

# Before box
sh = sl.shapes.add_shape(1, Inches(0.6), Inches(1.6),
                          Inches(5.15), Inches(5.25))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
sh.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); sh.line.width = Pt(1.5)
lbl(sl, Inches(0.85), Inches(1.75), Inches(4.6),
    "CODE REPO KNOWLEDGE LAYER", color=GRAY)
tb(sl, Inches(0.85), Inches(2.05), Inches(4.6), Inches(0.42),
   "Developer Tool", size=17, bold=True, color=DARK)
bullets(sl, Inches(0.85), Inches(2.55), Inches(4.6),
    ["Target: engineers in a team codebase",
     "Goal: faster onboarding & knowledge sharing",
     "Interface: VS Code Extension sidebar chat",
     "Requires: existing git & codebase context",
     "Needs: shared DB + server for team use"],
    size=11, gap=Inches(0.52), dot_col=GRAY, text_col=MID)

# Arrow
tb(sl, Inches(5.95), Inches(3.65), Inches(1.45), Inches(0.75),
   "->", size=40, bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(sl, Inches(5.95), Inches(4.4), Inches(1.45), Inches(0.3),
   "PIVOT", size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)

# After box
sh2 = sl.shapes.add_shape(1, Inches(7.55), Inches(1.6),
                           Inches(5.15), Inches(5.25))
sh2.fill.solid(); sh2.fill.fore_color.rgb = RED
sh2.line.fill.background(); sh2.line.width = 0
lbl(sl, Inches(7.8), Inches(1.75), Inches(4.6),
    "GIT ABSTRACTION TOOL", color=RGBColor(0xFF, 0xBB, 0xBB))
tb(sl, Inches(7.8), Inches(2.05), Inches(4.6), Inches(0.42),
   "Non-Technical User Tool", size=17, bold=True, color=WHITE)
bullets(sl, Inches(7.8), Inches(2.55), Inches(4.6),
    ["Target: vibe coders, designers, writers, students",
     "Goal: version control without knowing git",
     "Interface: Electron desktop app, plain language",
     "Requires: zero git knowledge",
     "Needs: local-only, no server, no cost"],
    size=11, gap=Inches(0.52), dot_col=WHITE, text_col=WHITE)

# ============================================================
# SLIDE 9 -- SECTION: Part 3
# ============================================================
sl = new_slide()
bg(sl, RED)
lbl(sl, Inches(1), Inches(2.7), Inches(10), "Part 3",
    color=RGBColor(0xFF, 0xCC, 0xCC))
h1(sl, Inches(1), Inches(3.05), Inches(11.3),
   "Git Abstraction Tool", size=38, color=WHITE)
tb(sl, Inches(1), Inches(4.2), Inches(9), Inches(0.6),
   "What we built -- and the 6 core AI features",
   size=17, color=RGBColor(0xFF, 0xDD, 0xDD))

# ============================================================
# SLIDE 10 -- GAT Overview
# ============================================================
sl = new_slide()
bg(sl, LIGHT)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "Git Abstraction Tool -- Overview")
h1(sl, Inches(0.6), Inches(0.62), Inches(12), "What It Is", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(1.4),
   "A cross-platform Electron desktop app that wraps git in a plain-language "
   "interface. No terminal. No git jargon. No learning curve.\n\n"
   "Built local-first: full version control works offline. AI features are "
   "optional enhancements layered on top, never requirements.",
   size=12, color=MID)

# Stats
for i, (num, desc) in enumerate([
        ("0", "git knowledge\nrequired"),
        ("6", "AI-powered\nfeatures"),
        ("2", "languages / EN+KO")]):
    sx = Inches(0.6 + i * 1.72)
    tb(sl, sx, Inches(3.15), Inches(1.6), Inches(0.62),
       num, size=34, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tb(sl, sx, Inches(3.75), Inches(1.6), Inches(0.5),
       desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Save Progress -- one-click local snapshot (git commit)",
     "Full File Visibility -- tracked, staged, untracked in one panel",
     "Safe Operations -- explicit confirmation for every risky action; "
     "no auto force-push",
     "Pull Preview -- inspect incoming changes before merging",
     "Bilingual UI -- English & Korean, beginner-friendly terminology",
     "Bring your own AI key -- OpenAI or Anthropic, stored encrypted locally"],
    size=12, gap=Inches(0.48))

# ============================================================
# SLIDE 11 -- Architecture
# ============================================================
sl = new_slide()
bg(sl, CREAM)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "Git Abstraction Tool -- Architecture")
h1(sl, Inches(0.6), Inches(0.62), Inches(12), "How It's Built", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

flow_items = [
    ("React UI",     "Renderer Process"),
    ("Electron IPC", "Secure Bridge"),
    ("Git Service",  "simple-git / Main"),
    ("SQLite DB",    "Local - No Server"),
    ("AI Layer",     "OpenAI / Anthropic"),
]
bw = Inches(2.05)
bh = Inches(1.45)
bg_gap = Inches(0.28)
total_w = len(flow_items) * bw + (len(flow_items) - 1) * bg_gap
sx = (W - total_w) / 2

for i, (name, sub) in enumerate(flow_items):
    bx = sx + i * (bw + bg_gap)
    by = Inches(2.0)
    sh = sl.shapes.add_shape(1, bx, by, bw, bh)
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); sh.line.width = Pt(1)
    tb(sl, bx, by + Inches(0.18), bw, Inches(0.4),
       name, size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(sl, bx, by + Inches(0.58), bw, Inches(0.32),
       sub, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        ax = bx + bw + Inches(0.04)
        tb(sl, ax, by + Inches(0.5), bg_gap - Inches(0.04), Inches(0.4),
           "->", size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

tech_items = [
    ("Frontend",
     "React + TypeScript, i18n (EN/KO), Tailwind -- renders in browser process"),
    ("Backend",
     "Electron main process handles all git ops, DB access, secret storage"),
    ("AI",
     "Modular provider adapters; OS-encrypted key storage; graceful degradation"),
]
cw = Inches(4.0)
for i, (h, b) in enumerate(tech_items):
    cx = Inches(0.6) + i * (cw + Inches(0.2))
    card(sl, cx, Inches(3.72), cw, Inches(1.0),
         h, b, hsize=11, bsize=10, accent=GRAY)

# ============================================================
# SLIDE 12 -- SECTION: AI Features
# ============================================================
sl = new_slide()
bg(sl, RED)
lbl(sl, Inches(1), Inches(2.1), Inches(11),
    "Git Abstraction Tool -- AI Features",
    color=RGBColor(0xFF, 0xCC, 0xCC))
h1(sl, Inches(1), Inches(2.5), Inches(11.3),
   "6 Core AI Components", size=38, color=WHITE)
tb(sl, Inches(1), Inches(3.75), Inches(11), Inches(0.55),
   "All optional  |  All privacy-preserving  |  All with graceful fallback",
   size=16, color=RGBColor(0xFF, 0xDD, 0xDD))

# ============================================================
# SLIDE 13 -- AI Feature 1: Auto Commit Message
# ============================================================
sl = new_slide()
bg(sl, LIGHT)
pill(sl, Inches(0.6), Inches(0.35), "AI Feature")
tb(sl, Inches(2.82), Inches(0.38), Inches(4), Inches(0.28),
   "AI Feature 1 of 6", size=8, bold=True, color=RED)
h1(sl, Inches(0.6), Inches(0.68), Inches(12),
   "Auto Commit Message", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(0.85),
   "When a user clicks Save Progress, the app reads the staged git diff and "
   "asks AI to produce one plain-language sentence -- no file paths, no git "
   "jargon.", size=12, color=MID)
quote(sl, Inches(0.6), Inches(2.58),  Inches(5.9),
      "Updated the homepage background color and fixed typos in the contact form.")
tb(sl, Inches(0.6), Inches(3.45), Inches(5.9), Inches(0.85),
   "The generated message is pre-filled for review. The internal AI summary "
   "is stored to power Weekly Report and Natural Language Undo.",
   size=11, color=GRAY)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Reads staged diff, file paths, and change types (new/modified/deleted)",
     "3-5 second timeout -- saving is never blocked by AI",
     "Per-project consent required before any diff is sent to AI",
     "Auto-Save mode: draft -> confirm -> commit in one flow",
     "Deterministic fallback message when AI is unavailable",
     "Stored AI summary powers the other three AI features"],
    size=12, gap=Inches(0.48))

# ============================================================
# SLIDE 14 -- AI Feature 2: Weekly Report
# ============================================================
sl = new_slide()
bg(sl, CREAM)
pill(sl, Inches(0.6), Inches(0.35), "AI Feature")
tb(sl, Inches(2.82), Inches(0.38), Inches(4), Inches(0.28),
   "AI Feature 2 of 6", size=8, bold=True, color=RED)
h1(sl, Inches(0.6), Inches(0.68), Inches(12), "Weekly Report", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(0.85),
   "At the end of any selected week, the app compiles git stats and stored "
   "AI summaries from each save into a readable narrative of what was "
   "accomplished.", size=12, color=MID)
card(sl, Inches(0.6), Inches(2.6), Inches(5.9), Inches(1.35),
     "Example Output",
     '"This week you redesigned the portfolio homepage, fixed spacing issues '
     'on mobile, and added a new contact form. 14 saves across 3 active days '
     '-- 847 lines added, 312 removed."', bsize=12)
tb(sl, Inches(0.6), Inches(4.1), Inches(5.9), Inches(0.6),
   "Useful for status updates, client logs, or personal portfolios -- "
   "no git knowledge needed.", size=11, color=GRAY)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Combines git stats with stored AI commit summaries for the week",
     "Shows total saves, files changed, lines added/removed, active days",
     "Cached by commit signature to avoid re-processing the same week",
     "Stats-only fallback (no AI narrative) when AI is not connected",
     "Distinguishes AI-summarized commits from message-only entries"],
    size=12, gap=Inches(0.52))

# ============================================================
# SLIDE 15 -- AI Feature 3: Natural Language Undo
# ============================================================
sl = new_slide()
bg(sl, WHITE)
pill(sl, Inches(0.6), Inches(0.35), "AI Feature")
tb(sl, Inches(2.82), Inches(0.38), Inches(4), Inches(0.28),
   "AI Feature 3 of 6", size=8, bold=True, color=RED)
h1(sl, Inches(0.6), Inches(0.68), Inches(12),
   "Natural Language Undo", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(0.65),
   "Instead of scrolling through commit hashes, users describe in plain "
   "language when they want to return to.", size=12, color=MID)
quote(sl, Inches(0.6), Inches(2.38), Inches(5.9),
      "Go back to before I changed the upload flow.")
quote(sl, Inches(0.6), Inches(3.2), Inches(5.9),
      "Restore to the state two days ago when the button was still red.")
tb(sl, Inches(0.6), Inches(4.05), Inches(5.9), Inches(0.75),
   "AI searches recent commits enriched with their AI summaries, shows a "
   "file-level preview, then waits for explicit confirmation before doing "
   "anything.", size=11, color=GRAY)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Searches up to 120 recent commits enriched with AI summaries",
     "Returns a primary match + 0-2 alternatives with confidence scores",
     "Shows which files would be restored and removed -- before applying",
     "User must explicitly confirm -- nothing applied automatically",
     "Fallback options when matched commit has no file-level delta"],
    size=12, gap=Inches(0.52))

# ============================================================
# SLIDE 16 -- AI Feature 4: Smart Conflict Resolver
# ============================================================
sl = new_slide()
bg(sl, CREAM)
pill(sl, Inches(0.6), Inches(0.35), "AI Feature")
tb(sl, Inches(2.82), Inches(0.38), Inches(4), Inches(0.28),
   "AI Feature 4 of 6", size=8, bold=True, color=RED)
h1(sl, Inches(0.6), Inches(0.68), Inches(12),
   "Smart Conflict Resolver", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(0.9),
   "When a pull or merge creates git conflicts, the app automatically "
   "surfaces a guided dialog -- no raw conflict markers exposed. "
   "Optional AI explains which version to keep and why.",
   size=12, color=MID)
card(sl, Inches(0.6), Inches(2.65), Inches(5.9), Inches(1.35),
     "AI Hint Example",
     '"Keep your version -- it has the new layout changes. The incoming '
     "version only adds a minor color fix that's already in yours.\"",
     bsize=12)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Dialog auto-opens when conflicts are detected after pull/merge",
     '"Keep Mine" / "Keep Theirs" per file -- no diff syntax required',
     "AI explains the tradeoff for each conflicted file on request",
     "Once resolved, AI can suggest the merge commit message",
     "Visual progress: resolved files show a checkmark badge",
     "Abort button always visible to cancel the entire merge safely"],
    size=12, gap=Inches(0.48))

# ============================================================
# SLIDE 17 -- AI Feature 5: File Insight
# ============================================================
sl = new_slide()
bg(sl, WHITE)
pill(sl, Inches(0.6), Inches(0.35), "AI Feature")
tb(sl, Inches(2.82), Inches(0.38), Inches(4), Inches(0.28),
   "AI Feature 5 of 6", size=8, bold=True, color=RED)
h1(sl, Inches(0.6), Inches(0.68), Inches(12),
   "File Insight", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(0.75),
   "Click any tracked file in the file panel and the app instantly "
   "generates a plain-language explanation of what that file does, "
   "how it works, and which related files connect to it.", size=12, color=MID)
card(sl, Inches(0.6), Inches(2.48), Inches(5.9), Inches(1.5),
     "Example Output",
     '"This is the main entry point of the React UI. It initializes the app, '
     'loads your saved preferences, and wires together the sidebar, file panel, '
     'and action panel. It coordinates the commit flow from file selection '
     'through to the save confirmation."',
     bsize=11)
tb(sl, Inches(0.6), Inches(4.12), Inches(5.9), Inches(0.6),
   "Explains any code file in plain language -- no terminal or code "
   "knowledge required.", size=11, color=GRAY)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Scores related files via co-commit history, directory proximity, "
     "and name similarity -- no hallucinated paths",
     "64 KB read limit; 12 KB content snippet sent to AI",
     "Binary files auto-detected by extension and null-byte scan and skipped",
     "Related file chips are clickable -- analyze connected files in sequence",
     "Requires AI connection; gracefully disabled when offline",
     "Request deduplication prevents stale results on rapid file switching"],
    size=12, gap=Inches(0.48))

# ============================================================
# SLIDE 18 -- AI Feature 6: Untracked File Review
# ============================================================
sl = new_slide()
bg(sl, LIGHT)
pill(sl, Inches(0.6), Inches(0.35), "AI Feature")
tb(sl, Inches(2.82), Inches(0.38), Inches(4), Inches(0.28),
   "AI Feature 6 of 6", size=8, bold=True, color=RED)
h1(sl, Inches(0.6), Inches(0.68), Inches(12),
   "Untracked File Review", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

tb(sl, Inches(0.6), Inches(1.6), Inches(5.9), Inches(0.75),
   "When untracked files exist and AI is connected, a 'Review untracked' "
   "button appears. The app classifies every untracked file: commit it "
   "to the repo, or delete it?", size=12, color=MID)
card(sl, Inches(0.6), Inches(2.48), Inches(5.9), Inches(1.5),
     "Example Recommendations",
     "Commit: src/Button.tsx (94%) -- 'Looks like new source code, likely "
     "belongs in version control.'\n"
     "Delete: node_modules/ (99%) -- 'Generated dependency folder, should "
     "never be committed to the repo.'",
     bsize=11)
tb(sl, Inches(0.6), Inches(4.12), Inches(5.9), Inches(0.6),
   "Stage all commit-recommended files in one click, or select delete-"
   "recommended files for instant removal.", size=11, color=GRAY)

bullets(sl, Inches(7.0), Inches(1.6), Inches(6.1),
    ["Two-tier decision: instant rule-based pass first, AI only for "
     "ambiguous cases -- fast and low token cost",
     "Rules auto-identify: build artifacts, caches, virtual envs, "
     ".env secrets, lock files, and source directories",
     "AI context: up to 100 files; up to 16 receive content snippets",
     "Adaptive timeout: 12s base + 220ms per file, capped at 45s",
     "Conservative fallback on timeout: unresolved files default to 'commit'",
     "Live UI updates as deletions complete; failed paths reported"],
    size=12, gap=Inches(0.48))

# ============================================================
# SLIDE 19 -- AI Summary Grid
# ============================================================
sl = new_slide()
bg(sl, CREAM)
lbl(sl, Inches(0.6), Inches(0.38), Inches(10),
    "Git Abstraction Tool -- AI Summary")
h1(sl, Inches(0.6), Inches(0.62), Inches(12),
   "6 Core AI Components at a Glance", size=28)
rule(sl, Inches(0.6), Inches(1.45), Inches(12.1))

ai_data = [
    ("Feature 1", "Auto Commit Message",
     "Turns staged diffs into a plain-language save description. Never "
     "blocks saving. Stored summaries power the other features."),
    ("Feature 2", "Weekly Report",
     "Generates a narrative of the week's work -- what was built, how many "
     "saves, across how many days. Stats-only fallback without AI."),
    ("Feature 3", "Natural Language Undo",
     "Describe when to return to in plain English. AI finds the commit, "
     "previews the file changes, and waits for your confirmation."),
    ("Feature 4", "Smart Conflict Resolver",
     "Guided dialog for merge conflicts -- file by file, plain language. "
     "AI explains which version to keep and drafts the merge message."),
    ("Feature 5", "File Insight",
     "Click any file for a plain-language explanation of what it does, "
     "how it works, and which related files connect to it."),
    ("Feature 6", "Untracked File Review",
     "Classifies untracked files as commit or delete. Rule-based pass "
     "first, AI handles ambiguous cases. Stage or delete in one click."),
]

cw = Inches(5.9)
ch = Inches(1.65)
gap = Inches(0.2)
for i, (num, head, body) in enumerate(ai_data):
    cx = Inches(0.6) + (i % 2) * (cw + Inches(0.43))
    cy = Inches(1.6) + (i // 2) * (ch + gap)
    sh = sl.shapes.add_shape(1, cx, cy, cw, ch)
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); sh.line.width = Pt(1)
    box(sl, cx, cy, cw, Inches(0.07), fill=RED)
    tb(sl, cx+Inches(0.18), cy+Inches(0.15), Inches(1.2), Inches(0.26),
       num.upper(), size=8, bold=True, color=RED)
    tb(sl, cx+Inches(0.18), cy+Inches(0.38), cw-Inches(0.36),
       Inches(0.35), head, size=13, bold=True, color=DARK)
    tb(sl, cx+Inches(0.18), cy+Inches(0.75), cw-Inches(0.36),
       Inches(0.78), body, size=10, color=MID)

# ============================================================
# SLIDE 18 -- CLOSING
# ============================================================
sl = new_slide()
bg(sl, DARK)
tb(sl, Inches(1), Inches(0.9), Inches(11.3), Inches(0.32),
   "SUMMARY", size=9, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
tb(sl, Inches(1), Inches(1.3), Inches(11.3), Inches(1.6),
   "Bringing version control to everyone --\nnot just developers.",
   size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider = sl.shapes.add_shape(1, Inches(5.6), Inches(3.1),
                               Inches(2.1), Pt(2))
divider.fill.solid(); divider.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
divider.line.fill.background(); divider.line.width = 0

tb(sl, Inches(0.8), Inches(3.25), Inches(5.3), Inches(0.55),
   "Code Repo Knowledge Layer",
   size=17, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(3.8), Inches(5.3), Inches(0.4),
   "Developer onboarding  |  VS Code extension",
   size=11, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

tb(sl, Inches(5.55), Inches(3.35), Inches(2.2), Inches(0.65),
   "->", size=36, bold=True, color=RED, align=PP_ALIGN.CENTER)

tb(sl, Inches(7.2), Inches(3.25), Inches(5.3), Inches(0.55),
   "Git Abstraction Tool",
   size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, Inches(7.2), Inches(3.8), Inches(5.3), Inches(0.4),
   "Version control for everyone  |  Electron desktop app",
   size=11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

tb(sl, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.35),
   "Honda x 99P Labs  |  Spring 2026",
   size=10, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
out = r"C:\Users\98sea\git_abstraction_tool\git-abstraction-tool\presentation.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides))
